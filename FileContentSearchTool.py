import os
import re
import chardet
import threading
from PyPDF2 import PdfReader
from pptx import Presentation
import customtkinter as ctk
from tkinter import filedialog, messagebox, Listbox, ttk, BooleanVar, END
import xml.etree.ElementTree as ET

class FolderSearchApp(ctk.CTk):
    def __init__(self):
        super().__init__()

        self.title("File Content Search Tool Ver. 2.0.0")
        self.geometry("1024x800")

        self.file_path_mapping = {}
        self.folder_list = []
        self.lock = threading.Lock()

        self.load_folders()

        self.include_subdirectories = BooleanVar(value=False)

        self.create_widgets()

    def create_widgets(self):
        self.folder_list_label = ctk.CTkLabel(self, text="路徑設定", font=("Microsoft JhengHei", 14))
        self.folder_list_label.pack(pady=5)

        folder_frame = ctk.CTkFrame(self)
        folder_frame.pack(fill=ctk.BOTH, expand=True, padx=int(self.winfo_width() * 0.025), pady=5)

        self.folder_listbox = Listbox(folder_frame, height=1)
        self.folder_listbox.pack(side=ctk.LEFT, fill=ctk.BOTH, expand=True)
        self.folder_scrollbar = ctk.CTkScrollbar(folder_frame, orientation='vertical', command=self.folder_listbox.yview)
        self.folder_scrollbar.pack(side=ctk.RIGHT, fill=ctk.Y)
        self.folder_listbox.config(yscrollcommand=self.folder_scrollbar.set)

        self.update_folder_listbox()

        self.button_frame = ctk.CTkFrame(self, fg_color='transparent')
        self.button_frame.pack(fill=ctk.X, pady=5)

        self.add_folder_button = ctk.CTkButton(self.button_frame, text="Add Folder", command=self.add_folder)
        self.add_folder_button.pack(side=ctk.LEFT, padx=10, expand=True, anchor='center')

        self.delete_folder_button = ctk.CTkButton(self.button_frame, text="Delete Folder", command=self.delete_folder)
        self.delete_folder_button.pack(side=ctk.LEFT, padx=10, expand=True, anchor='center')

        self.search_results_label = ctk.CTkLabel(self, text="搜尋功能 (結果雙點即可自動開啟)", font=("Microsoft JhengHei", 14))
        self.search_results_label.pack(pady=(20, 5))

        self.search_frame = ctk.CTkFrame(self, fg_color='transparent')
        self.search_frame.pack(fill=ctk.X, pady=5)

        self.search_entry = ctk.CTkEntry(self.search_frame, placeholder_text="Enter keyword")
        self.search_entry.pack(side=ctk.LEFT, expand=True, fill=ctk.X, padx=10)

        self.include_subdirs_checkbutton = ctk.CTkCheckBox(self.search_frame, text="包含子目錄", variable=self.include_subdirectories)
        self.include_subdirs_checkbutton.pack(side=ctk.LEFT, padx=10)

        self.search_button = ctk.CTkButton(self.search_frame, text="Search", command=self.search)
        self.search_button.pack(side=ctk.LEFT, padx=10)

        # Create a Treeview for search results
        result_frame = ctk.CTkFrame(self)
        result_frame.pack(fill=ctk.BOTH, expand=True, padx=int(self.winfo_width() * 0.025), pady=5)

        self.search_results = ttk.Treeview(result_frame, columns=("FileName", "Type", "Content", "Path"), show="headings")
        self.search_results.heading("FileName", text="FileName")
        self.search_results.heading("Type", text="Type")
        self.search_results.heading("Content", text="Content")
        self.search_results.heading("Path", text="Path")
        self.search_results.column("FileName", width=350, anchor="w")
        self.search_results.column("Type", width=50, anchor="w")
        self.search_results.column("Content", width=500, anchor="w")
        self.search_results.column("Path", width=50, anchor="w")
        self.search_results.pack(side=ctk.LEFT, fill=ctk.BOTH, expand=True)

        self.search_scrollbar = ctk.CTkScrollbar(result_frame, orientation='vertical', command=self.search_results.yview)
        self.search_scrollbar.pack(side=ctk.RIGHT, fill=ctk.Y)

        self.search_results.config(yscrollcommand=self.search_scrollbar.set)
        self.search_results.bind("<Double-Button-1>", self.open_file)

        self.progress_bar = ctk.CTkProgressBar(self, height=20, progress_color="#90EE90")
        self.progress_bar.pack(fill=ctk.X, padx=self.winfo_width() * 0.025, pady=(5, 5))
        self.progress_bar.set(0)

    def add_folder(self):
        folder_path = filedialog.askdirectory()
        if folder_path:
            if folder_path not in self.folder_list:
                self.folder_list.append(folder_path)
                self.update_folder_listbox()
                self.save_folders()

    def delete_folder(self):
        selected_indices = self.folder_listbox.curselection()
        if selected_indices:
            selected_folder = self.folder_listbox.get(selected_indices[0])
            if selected_folder in self.folder_list:
                self.folder_list.remove(selected_folder)
            self.folder_listbox.delete(selected_indices[0])
            self.save_folders()

    def update_folder_listbox(self):
        self.folder_listbox.delete(0, END)
        for folder in self.folder_list:
            self.folder_listbox.insert(0, folder)

    def open_file(self, event):
        selected_item = self.search_results.focus()
        if selected_item:
            selected_file_path = self.search_results.item(selected_item)["values"][3]
            if selected_file_path:
                os.startfile(selected_file_path)

    def save_folders(self):
        root = ET.Element("folders")
        for folder in self.folder_list:
            ET.SubElement(root, "folder").text = folder

        tree = ET.ElementTree(root)
        tree.write("folders.xml")

    def load_folders(self):
        try:
            tree = ET.parse("folders.xml")
            root = tree.getroot()
            self.folder_list = [folder.text for folder in root.findall("folder")]
        except FileNotFoundError:
            self.folder_list = []

    def update_progress_bar(self, progress_increment):
        with self.lock:
            new_progress = self.progress_bar.get() + progress_increment
            self.progress_bar.set(new_progress)
            self.update()

    def search(self):
        keyword = self.search_entry.get()
        if not keyword:
            messagebox.showwarning("Warning", "請輸入關鍵字進行搜索。")
            return

        keyword_lower = keyword.lower()

        self.search_results.delete(*self.search_results.get_children())

        self.progress_bar.set(0)
        self.update()

        total_files = self.calculate_total_files()
        if total_files >= 100:
            user_confirmation = messagebox.askyesno("Warning", f"檔案數量共 {total_files} 個，由於數量過多可能會搜尋很久，確定要執行嗎?")
            if not user_confirmation:
                return
 
        if total_files == 0:
            messagebox.showwarning("Warning", "路徑內沒有檔案。")
            return

        progress_increment = 1.0 / total_files

        for folder in self.folder_list:
            for root, dirs, files in os.walk(folder, topdown=True):
                if not self.include_subdirectories.get():
                    dirs.clear()
                for file in files:
                    file_path = os.path.join(root, file)
                    if file.endswith(".pdf"):
                        self.search_in_file(file_path, keyword_lower, progress_increment, file_type="pdf")
                    elif file.endswith(".pptx"):
                        self.search_in_file(file_path, keyword_lower, progress_increment, file_type="pptx")
                    else:
                        self.search_in_file(file_path, keyword_lower, progress_increment, file_type="txt")

        if not self.search_results.get_children():
            messagebox.showinfo("Information", "搜索完畢，但結果是空的。")

    def calculate_total_files(self):
        total_files = 0
        for folder in self.folder_list:
            for root, dirs, files in os.walk(folder, topdown=True):
                if not self.include_subdirectories.get():
                    dirs.clear()
                for file in files:
                    # if file.endswith(".pdf") or file.endswith(".pptx") or file.endswith(".txt") or file.endswith(".ini"):
                    total_files += 1
        return total_files

    def search_in_file(self, file_path, keyword_lower, progress_increment, file_type):
        try:
            self.update_progress_bar(progress_increment)
            if file_type == "pdf":
                with open(file_path, "rb") as file:
                    content = PdfReader(file)
                    self.search_in_pdf(content, keyword_lower, file_path)
            elif file_type == "pptx":
                with open(file_path, "rb") as file:
                    content = Presentation(file)
                    self.search_in_pptx(content, keyword_lower, file_path)
            else:
                self.search_in_text_file(file_path, keyword_lower, progress_increment, file_type)
        except Exception as e:
            result_text = f"處理文件時出現錯誤「{e}」"
            self.search_results.insert("", "end", values=("Error", file_path, result_text))

    def search_in_pdf(self, content, keyword_lower, file_path):
        file_exists = False
        for page_num, page in enumerate(content.pages):
            original_text = page.extract_text()
            text_lower = original_text.lower()
            if keyword_lower in text_lower:
                lines = original_text.split("\n")
                for line in lines:
                    line_lower = line.lower()
                    if keyword_lower in line_lower:
                        file_name = os.path.basename(file_path)
                        file_path = re.sub(r"/", r"\\", file_path)
                        self.search_results.insert("", "end", values=(file_name, "PDF", line, file_path))
                        file_exists = True
                        break
            if file_exists:
                break

    def search_in_pptx(self, content, keyword_lower, file_path):
        file_exists = False
        for slide_num, slide in enumerate(content.slides):
            original_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    original_text.append(shape.text_frame.text)
            slide_text = " ".join(original_text)
            text_lower = slide_text.lower()
            if keyword_lower in text_lower:
                lines = slide_text.split("\n")
                for line in lines:
                    line_lower = line.lower()
                    if keyword_lower in line_lower:
                        file_name = os.path.basename(file_path)
                        file_path = re.sub(r"/", r"\\", file_path)
                        self.search_results.insert("", "end", values=(file_name, "PPT", line, file_path))
                        file_exists = True
                        break
            if file_exists:
                break

    def search_in_text_file(self, file_path, keyword_lower, progress_increment, file_type):
        try:
            self.update_progress_bar(progress_increment)
            
            with open(file_path, "rb") as file:
                raw_data = file.read(100)
                result = chardet.detect(raw_data)
                encoding = result['encoding']

            with open(file_path, "r", encoding=encoding, errors='ignore') as file:
                for line in file:
                    line_lower = line.lower()
                    if keyword_lower in line_lower:
                        file_name = os.path.basename(file_path)
                        file_path = re.sub(r"/", r"\\", file_path)
                        self.search_results.insert("", "end", values=(file_name, file_type.upper(), line.strip(), file_path))
                        break
        except Exception as e:
            result_text = f"處理文件時出現錯誤「{e}」"
            self.search_results.insert("", "end", values=("Error", file_path, result_text))
            print(result_text)

if __name__ == "__main__":
    app = FolderSearchApp()
    app.mainloop()
